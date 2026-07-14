#!/usr/bin/env python3
"""Autonomous .pptx sanity check — no LibreOffice, no admin.

    python validate_pptx.py out.pptx

Checks the failure modes that actually break template-derived decks:
  * every part is well-formed XML
  * every slide has a Content-Types Override
  * every sldId r:id resolves to a relationship whose target slide exists
  * each slide's internal relationship targets exist
  * python-pptx can open the deck and iterate its slides (deep smoke test)

Exits non-zero and prints each problem if anything fails. This is a pragmatic
check, not full XSD schema validation; if python-pptx opens the deck and the
relationship graph is intact, PowerPoint opens it in practice.
"""
import os
import sys
import zipfile
from defusedxml.minidom import parseString

SLIDE_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"


def main():
    if len(sys.argv) < 2:
        sys.exit("usage: python validate_pptx.py deck.pptx")
    # accept and ignore extra flags (e.g. --original) for CLI parity
    path = next((a for a in sys.argv[1:] if not a.startswith("-")), None)
    if not path or not os.path.exists(path):
        sys.exit(f"error: file not found: {path}")

    problems = []
    with zipfile.ZipFile(path) as z:
        names = set(z.namelist())

        # 1. well-formed XML for every xml/rels part
        for n in names:
            if n.endswith(".xml") or n.endswith(".rels"):
                try:
                    parseString(z.read(n))
                except Exception as e:
                    problems.append(f"malformed XML in {n}: {e}")

        # 2. content-types Override for every slide
        if "[Content_Types].xml" not in names:
            problems.append("missing [Content_Types].xml")
        else:
            ct = parseString(z.read("[Content_Types].xml"))
            overrides = {o.getAttribute("PartName")
                         for o in ct.getElementsByTagName("Override")}
            for n in names:
                if n.startswith("ppt/slides/slide") and n.endswith(".xml"):
                    if "/" + n not in overrides:
                        problems.append(f"no Content-Types Override for /{n}")

        # 3. sldId r:id -> relationship -> existing slide part
        pres_rels = "ppt/_rels/presentation.xml.rels"
        rid_target = {}
        if pres_rels in names:
            pr = parseString(z.read(pres_rels))
            for r in pr.getElementsByTagName("Relationship"):
                rid_target[r.getAttribute("Id")] = r.getAttribute("Target")
        if "ppt/presentation.xml" in names:
            pres = parseString(z.read("ppt/presentation.xml"))
            for s in pres.getElementsByTagName("p:sldId"):
                rid = s.getAttribute("r:id")
                if rid not in rid_target:
                    problems.append(f"sldId {rid} has no relationship in presentation.xml.rels")
                    continue
                tgt = "ppt/" + rid_target[rid].replace("../", "")
                if tgt not in names:
                    problems.append(f"sldId {rid} -> {tgt} but that part is missing")

        # 4. each slide's internal relationship targets exist
        for n in list(names):
            if n.startswith("ppt/slides/_rels/") and n.endswith(".rels"):
                base = "ppt/slides/"
                rr = parseString(z.read(n))
                for r in rr.getElementsByTagName("Relationship"):
                    if r.getAttribute("TargetMode") == "External":
                        continue
                    tgt = os.path.normpath(base + r.getAttribute("Target")).replace("\\", "/")
                    if tgt not in names:
                        problems.append(f"{n}: target missing -> {tgt}")

    # 5. deep smoke test via python-pptx
    try:
        from pptx import Presentation
        prs = Presentation(path)
        _ = len(prs.slides._sldIdLst)  # force parse
        for _slide in prs.slides:
            _ = _slide.shapes
    except ImportError:
        problems.append("python-pptx not installed (pip install --user python-pptx) — "
                        "skipped deep open test")
    except Exception as e:
        problems.append(f"python-pptx could not open the deck: {e}")

    if problems:
        print("VALIDATION FAILED:")
        for p in problems:
            print("  -", p)
        sys.exit(1)
    print("All validations PASSED!")


if __name__ == "__main__":
    main()
